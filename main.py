from pptx import Presentation
import argparse
import sys, time

def generator(args):

    in_file, out_file = args.infile, args.outfile
    prs = Presentation(in_file)
    horizontal, column, images = get_setting(prs, args)

    # 1. 가로 세로를 순서대로 정렬하는 코드
    horizontal = sorted(horizontal, key=lambda x: int(x.top))
    column = sorted(column, key=lambda x: int(x.left))

    # 2. 가로의 개수만큼씩 나누는 코드.
    images = sorted(images, key=lambda x: int(x.top))
    temp = [[] for _ in range(len(column))]

    for i in range(len(column)):
        temp[i] = images[i*len(horizontal): (i+1)*len(horizontal)]
        temp[i] = sorted(temp[i], key=lambda x: int(x.left))
    images = temp

    # 3. align
    for i in range(len(column)):
        for j in range(len(horizontal)):
            images[i][j].top = column[i].top
            images[i][j].left = horizontal[j].left

    prs.save(out_file)

def get_setting(prs, args):
    shapes = prs.slides[args.where-1].shapes
    horizontal = []
    column = []
    images = []
    for shape in shapes:
        if int(shape.left) < 0:
            column.append(shape)
        elif int(shape.top) < 0:
            horizontal.append(shape)
        else:
            images.append(shape)

    # number check
    if (len(images) != len(column) * len(horizontal)):
        print("The number of images is not same with grid shape")
        time.sleep(3)
        sys.exit()

    return horizontal, column, images

def explain():
    print("################################")
    print("1. 파일 개수가 grid 형식과 다르면 작동하지 않음.")
    print("2. input ppt의 첫번째 slide에서 작동하는 것이 default. 바꾸고 싶으면 --where 인자를 사용할 것. (index는 1부터 시작)")
    print("################################")

if __name__ == "__main__":
    explain()

    parser = argparse.ArgumentParser()
    parser.add_argument('--infile', type=str, help='input file path')
    parser.add_argument('--outfile', type=str, help='output file path')
    parser.add_argument('--where', type=int, default=1,  help='output file path')
    args = parser.parse_args()

    generator(args)

